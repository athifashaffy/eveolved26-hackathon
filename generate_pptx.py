"""Generate PowerPoint from slides content."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Colors
BG = RGBColor(0x0a, 0x0a, 0x1a)
ACCENT = RGBColor(0x00, 0xd4, 0xaa)
ACCENT2 = RGBColor(0x6c, 0x5c, 0xe7)
WHITE = RGBColor(0xff, 0xff, 0xff)
LIGHT = RGBColor(0xe0, 0xe0, 0xe0)
DIM = RGBColor(0x88, 0x88, 0x88)
WARN = RGBColor(0xff, 0x6b, 0x6b)
BOX_BG = RGBColor(0x16, 0x16, 0x2a)
BOX_ACCENT_BG = RGBColor(0x0d, 0x20, 0x1c)


def set_slide_bg(slide, color=BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, left, top, width, height, text, font_size=18,
             color=LIGHT, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_box(slide, left, top, width, height, fill_color=BOX_BG):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top),
        Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = RGBColor(0x33, 0x33, 0x44)
    shape.line.width = Pt(1)
    return shape


def add_stat(slide, left, top, number, label):
    add_text(slide, left, top, 2.8, 0.8, number, font_size=48, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, left, top + 0.7, 2.8, 0.4, label, font_size=12, color=DIM, alignment=PP_ALIGN.CENTER)


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ============ SLIDE 1: Title ============
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide)
add_text(slide, 1, 1.5, 11, 1.2, "AutoBiomarker", font_size=54, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text(slide, 1, 2.8, 11, 0.8, "Autonomous Discovery of Depression Early Warning Signals\nfrom Wearable Data", font_size=22, color=DIM, alignment=PP_ALIGN.CENTER)
add_text(slide, 1, 4.2, 11, 0.5, "Evolved26 Toronto — March 2026", font_size=14, color=DIM, alignment=PP_ALIGN.CENTER)
add_text(slide, 1, 5.2, 11, 0.5, "Nebius Token Factory  •  Karpathy Autoresearch  •  Temporal Dynamics Theory", font_size=12, color=DIM, alignment=PP_ALIGN.CENTER)

# ============ SLIDE 2: The Problem ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text(slide, 0.5, 0.3, 12, 0.7, "The Problem", font_size=36, color=ACCENT, bold=True)
add_stat(slide, 1.5, 1.2, "280M", "people with depression globally")
add_stat(slide, 5.2, 1.2, "50%", "drop out of treatment")
add_stat(slide, 8.9, 1.2, "75%", "of depressed patients\nhave sleep problems")
add_box(slide, 1, 3.2, 11, 1.8)
add_text(slide, 1.3, 3.4, 10.4, 0.6, "Depression is gradual — patients don't \"fall into\" it, they drift.", font_size=18, color=LIGHT)
add_text(slide, 1.3, 3.9, 10.4, 0.6, "By the time they report worsening, they've often dropped out.", font_size=18, color=WARN)
add_text(slide, 1.3, 4.4, 10.4, 0.6, "But sleep and physiology destabilize before patients notice — and sleep is treatable.", font_size=18, color=ACCENT)

# ============ SLIDE 3: Insight ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text(slide, 0.5, 0.3, 12, 0.7, "The Insight: Destabilization Before Decline", font_size=32, color=ACCENT, bold=True)
add_box(slide, 0.8, 1.2, 11.5, 1.2, fill_color=RGBColor(0x1a, 0x15, 0x30))
add_text(slide, 1.1, 1.3, 11, 1.0, "Depression worsens gradually. But as physiological regulation weakens,\nthe body shows measurable signs: sleep becomes irregular, HRV becomes erratic, recovery slows.", font_size=17, color=LIGHT)
add_box(slide, 0.8, 2.8, 5.3, 2.2)
add_text(slide, 1.1, 2.9, 4.8, 0.4, "Known from physics", font_size=18, color=ACCENT2, bold=True)
add_text(slide, 1.1, 3.3, 4.8, 1.2, "Systems under stress show rising variability\nand sluggish recovery — \"critical slowing down\"\n\nProven in: Ecology • Climate • Finance • Mood (PNAS 2014)", font_size=14, color=LIGHT)
add_box(slide, 6.5, 2.8, 5.3, 2.2, fill_color=BOX_ACCENT_BG)
add_text(slide, 6.8, 2.9, 4.8, 0.4, "Our hypothesis", font_size=18, color=ACCENT2, bold=True)
add_text(slide, 6.8, 3.3, 4.8, 1.2, "The same temporal patterns in wearable data\ncan flag when a patient's regulation is destabilizing\n\nNot a sudden tipping point — a gradual loss of stability", font_size=14, color=LIGHT)

# ============ SLIDE 4: Our Approach ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text(slide, 0.5, 0.3, 12, 0.7, "What Everyone Measures vs. What We Measure", font_size=30, color=ACCENT, bold=True)
add_box(slide, 0.8, 1.3, 5.3, 2.0)
add_text(slide, 1.1, 1.4, 4.8, 0.4, "Standard Approach", font_size=18, color=DIM, bold=True)
add_text(slide, 1.1, 1.8, 4.8, 1.2, "How low is your HRV?\n\nSingle timepoint, cross-sectional\nSmall effect sizes (d ≈ 0.3)\nDepression-specific? No — stress marker", font_size=14, color=DIM)
add_box(slide, 6.5, 1.3, 5.3, 2.0, fill_color=BOX_ACCENT_BG)
add_text(slide, 6.8, 1.4, 4.8, 0.4, "Our Approach", font_size=18, color=ACCENT, bold=True)
add_text(slide, 6.8, 1.8, 4.8, 1.2, "How unstable is your physiology over time?\n\nTemporal dynamics: autocorrelation, variance, CV\nDay-to-day patterns across 28 days\nTransdiagnostic — tracks stress regulation", font_size=14, color=LIGHT)
add_box(slide, 0.8, 3.7, 11, 1.2)
add_text(slide, 1.1, 3.8, 10.4, 1.0, "HRV and sleep are transdiagnostic stress markers (not depression-specific).\nThis is a strength: our system monitors overall physiological regulation,\nwith sleep as a directly treatable intervention target.", font_size=14, color=LIGHT)

# ============ SLIDE 5: Autoresearch Loop ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text(slide, 0.5, 0.3, 12, 0.7, "The Autoresearch Loop", font_size=36, color=ACCENT, bold=True)
add_text(slide, 0.5, 0.9, 12, 0.4, "Adapted from Karpathy's autoresearch for biomarker discovery", font_size=13, color=DIM)
# Flow boxes
for i, (label, sub) in enumerate([
    ("Hypothesis\nGenerator", "Llama 3.3 70B\non Nebius"),
    ("Feature\nExtractor", "187+ temporal\nfeatures"),
    ("Statistical\nEvaluator", "LOSO-CV,\nAUC, FDR"),
    ("KEEP /\nDISCARD", ""),
]):
    x = 1.2 + i * 3.0
    c = BOX_ACCENT_BG if i == 3 else RGBColor(0x1a, 0x15, 0x30)
    add_box(slide, x, 1.6, 2.3, 1.3, fill_color=c)
    add_text(slide, x + 0.1, 1.7, 2.1, 0.6, label, font_size=15, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    if sub:
        add_text(slide, x + 0.1, 2.3, 2.1, 0.5, sub, font_size=11, color=DIM, alignment=PP_ALIGN.CENTER)
    if i < 3:
        add_text(slide, x + 2.3, 1.9, 0.7, 0.5, "→", font_size=28, color=ACCENT, alignment=PP_ALIGN.CENTER)
add_text(slide, 1, 3.3, 11, 0.4, "← ← ← ← ← ← ← ← ← Results feed back into hypothesis generation ← ← ← ← ← ← ← ← ←", font_size=12, color=ACCENT, alignment=PP_ALIGN.CENTER)
add_box(slide, 1, 3.9, 11, 0.8)
add_text(slide, 1.3, 4.0, 10.4, 0.6, "Set it up before sleep. Wake up to 5,000 hypotheses tested with full statistical rigor.", font_size=16, color=LIGHT, alignment=PP_ALIGN.CENTER)

# ============ SLIDE 6: Results ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text(slide, 0.5, 0.3, 12, 0.7, "Results: Baigutanova Dataset", font_size=36, color=ACCENT, bold=True)
add_stat(slide, 1.5, 1.2, "5,000", "feature-tests run autonomously")
add_stat(slide, 5.2, 1.2, "3", "significant after FDR\n(p_adj < 0.05)")
add_stat(slide, 8.9, 1.2, "0.76", "best interaction AUC")
add_box(slide, 1, 3.2, 11, 2.0, fill_color=BOX_ACCENT_BG)
add_text(slide, 1.3, 3.3, 10.4, 0.5, "Top: LF/HF 21d slope x sleep quality — AUC 0.76, d=1.06, p_adj=0.02", font_size=18, color=ACCENT, bold=True)
add_text(slide, 1.3, 3.8, 10.4, 0.5, "Also: RMSSD slope x sleep dur autocorr (d=1.11, p_adj=0.006), RMSSD slope x LF/HF slope (d=0.92, p_adj=0.004)", font_size=14, color=LIGHT)
add_text(slide, 1.3, 4.3, 10.4, 0.5, "All significant features are HRV x sleep interactions — multi-system destabilization is key", font_size=14, color=LIGHT)
add_text(slide, 1.3, 4.8, 10.4, 0.3, "50 LLM-guided rounds (271 tests), per-round BH-FDR at n=49. Total LLM cost: $0.22.", font_size=11, color=DIM)

# ============ SLIDE 7: Trained Model ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text(slide, 0.5, 0.3, 12, 0.7, "Trained Prediction Model", font_size=34, color=ACCENT, bold=True)
add_stat(slide, 1.5, 1.1, "0.76", "AUC (best interaction)")
add_stat(slide, 5.2, 1.1, "0.69", "AUC (LOSO-CV model)")
add_stat(slide, 8.9, 1.1, "49", "subjects (LOSO-CV)")
add_box(slide, 0.8, 2.8, 5.5, 2.4)
add_text(slide, 1.1, 2.9, 5, 0.4, "Architecture", font_size=18, color=ACCENT2, bold=True)
add_text(slide, 1.1, 3.3, 5, 1.6, "• Feature selection INSIDE each CV fold (no leakage)\n• Leave-one-subject-out cross-validation\n• 732 candidate features → 50 selected per fold\n• MLP Neural Net + autoresearch features", font_size=13, color=LIGHT)
add_box(slide, 6.7, 2.8, 5.5, 2.4, fill_color=BOX_ACCENT_BG)
add_text(slide, 7.0, 2.9, 5, 0.4, "Significant Interactions (FDR<0.05)", font_size=18, color=ACCENT2, bold=True)
add_text(slide, 7.0, 3.3, 5, 1.6, "• LF/HF slope × sleep quality (AUC 0.76, d=1.06)\n• RMSSD slope × sleep dur autocorr (d=1.11)\n• RMSSD slope × LF/HF slope (d=0.92)\n• LF/HF slope × sleep latency (AUC 0.70)\n• All cross HRV × sleep domains", font_size=13, color=LIGHT)
add_text(slide, 0.8, 5.4, 11.5, 0.4, "HRV × sleep interactions outperform single features — multi-system destabilization is the signal.", font_size=13, color=ACCENT, alignment=PP_ALIGN.CENTER)

# ============ SLIDE 8: Self-Improving Loop ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text(slide, 0.5, 0.3, 12, 0.7, "Self-Improving Discovery Engine", font_size=34, color=ACCENT, bold=True)
# Flow
for i, (label, sub) in enumerate([
    ("Round 1", "Broad sweep"), ("5,000 tested", "137 promising"),
    ("Feed back\nto LLM", "Top findings + nulls"), ("50 rounds", "3 significant!"),
]):
    x = 1.2 + i * 3.0
    c = BOX_ACCENT_BG if i >= 2 else RGBColor(0x1a, 0x15, 0x30)
    add_box(slide, x, 1.2, 2.3, 1.1, fill_color=c)
    add_text(slide, x + 0.1, 1.25, 2.1, 0.5, label, font_size=15, color=WHITE if i != 3 else ACCENT, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, x + 0.1, 1.75, 2.1, 0.4, sub, font_size=11, color=DIM if i != 3 else ACCENT, alignment=PP_ALIGN.CENTER)
    if i < 3:
        add_text(slide, x + 2.3, 1.4, 0.7, 0.5, "→", font_size=28, color=ACCENT, alignment=PP_ALIGN.CENTER)
add_box(slide, 0.8, 2.8, 5.3, 1.5)
add_text(slide, 1.1, 2.9, 4.8, 0.3, "Round 1 (broad)", font_size=16, color=DIM, bold=True)
add_text(slide, 1.1, 3.2, 4.8, 0.9, "5,000 hypotheses • 137 promising pre-FDR\n0 after strict BH-FDR across 5,000 tests", font_size=14, color=LIGHT)
add_box(slide, 6.5, 2.8, 5.3, 1.5, fill_color=BOX_ACCENT_BG)
add_text(slide, 6.8, 2.9, 4.8, 0.3, "Rounds 2-50 (LLM-guided)", font_size=16, color=ACCENT, bold=True)
add_text(slide, 6.8, 3.2, 4.8, 0.9, "271 targeted hypotheses • 3 survived FDR\nBest: lfhf_slope × sleep_qual (AUC 0.76)\nTotal cost: $0.22 on Nebius", font_size=14, color=LIGHT)
add_text(slide, 0.8, 4.6, 11.5, 0.4, "Small rounds (8 tests) keep FDR mild. LLM adapts each round. $0.22 total for 50 rounds of discovery.", font_size=12, color=DIM, alignment=PP_ALIGN.CENTER)

# ============ SLIDE 9: Methods ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text(slide, 0.5, 0.3, 12, 0.7, "Methods: Model & Evaluation", font_size=34, color=ACCENT, bold=True)
add_box(slide, 0.8, 1.2, 5.3, 2.8, fill_color=BOX_ACCENT_BG)
add_text(slide, 1.1, 1.3, 4.8, 0.3, "Data Split", font_size=16, color=ACCENT, bold=True)
add_text(slide, 1.1, 1.7, 4.8, 2.0, "49 subjects (Baigutanova et al. 2025)\nStratified by PHQ-9 status\n\nDiscovery: 33 subjects (70%)\nValidation: 16 subjects (30%)\n\nClass ratio preserved in both sets", font_size=13, color=LIGHT)
add_box(slide, 6.5, 1.2, 5.3, 2.8)
add_text(slide, 6.8, 1.3, 4.8, 0.3, "Classifier Pipeline", font_size=16, color=ACCENT2, bold=True)
add_text(slide, 6.8, 1.7, 4.8, 2.0, "1. 187 temporal features (autocorr, CV, slope, std)\n   over 3/7/14/21-day windows\n2. HRV x sleep cross-domain interactions\n3. Per-subject z-scored composites\n4. LOSO-CV — feature selection INSIDE each fold\n5. Logistic Regression with top-10 features", font_size=13, color=LIGHT)
add_box(slide, 0.8, 4.3, 11.5, 0.8)
add_text(slide, 1.1, 4.4, 11, 0.6, "Key: Feature selection inside CV folds prevents data leakage. Subject-level stats prevent pseudoreplication.", font_size=14, color=ACCENT, alignment=PP_ALIGN.CENTER)

# ============ SLIDE 10: Validation ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text(slide, 0.5, 0.3, 12, 0.7, "Validation: Is It Overfitted?", font_size=34, color=ACCENT, bold=True)
add_text(slide, 0.5, 0.9, 12, 0.3, "We ran 4 independent validation tests to verify our findings are real.", font_size=13, color=DIM)
add_box(slide, 0.8, 1.4, 5.7, 3.0, fill_color=BOX_ACCENT_BG)
add_text(slide, 1.1, 1.5, 5.2, 0.3, "Permutation Test (gold standard)", font_size=16, color=ACCENT, bold=True)
add_text(slide, 1.1, 1.9, 5.2, 2.2, "Shuffled depression labels 100 times, re-ran analysis\n\nRMSSD slope x sleep dur:  real d=1.11, null=0.26  →  0/100 beat\nRMSSD slope x LF/HF:      real d=0.92, null=0.22  →  0/100 beat\nLF/HF slope x sleep qual:  real d=1.06, null=0.30  →  3/100 beat\n\n0 out of 20 full shuffled simulations produced\nANY significant finding. Real data: 3.", font_size=12, color=LIGHT)
add_box(slide, 6.8, 1.4, 5.7, 3.0)
add_text(slide, 7.1, 1.5, 5.2, 0.3, "Hold-Out Validation", font_size=16, color=ACCENT2, bold=True)
add_text(slide, 7.1, 1.9, 5.2, 2.2, "33 discovery / 16 validation (stratified split)\n\nRMSSD x sleep dur:   disc d=1.09 → val d=1.08  ✓ Consistent\nRMSSD x LF/HF:       disc d=0.55 → val d=1.95  ✓ Consistent\nLF/HF x sleep qual:   disc d=0.89 → val d=1.54  ✓ Consistent\n\nAll 3 effect directions consistent.\n2/3 reach p<0.05 in held-out set.", font_size=12, color=LIGHT)
add_box(slide, 0.8, 4.7, 11.5, 0.8)
add_text(slide, 1.1, 4.8, 11, 0.6, "Verdict: Not overfitted. Real effects 3-4x larger than null. Effects replicate in unseen subjects.", font_size=15, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)

# ============ SLIDE 11: Clinical Utility ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text(slide, 0.5, 0.3, 12, 0.7, "Why This Matters Clinically", font_size=36, color=ACCENT, bold=True)
add_box(slide, 0.8, 1.1, 11.5, 1.3, fill_color=RGBColor(0x1a, 0x15, 0x30))
add_text(slide, 1.1, 1.2, 11, 0.8, "\"The good thing with sleep problems is: if you treat them it hits two birds\nwith one stone — a bothering problem that also triggers other symptoms.\"", font_size=17, color=LIGHT)
add_text(slide, 1.1, 2.0, 11, 0.3, "— Prof. Dr. Steffen Moritz, Uni of Hamburg", font_size=12, color=DIM)
add_box(slide, 0.8, 2.8, 5.3, 2.0)
add_text(slide, 1.1, 2.9, 4.8, 0.3, "The Problem with Biomarkers", font_size=16, color=ACCENT2, bold=True)
add_text(slide, 1.1, 3.2, 4.8, 1.2, "• Predictive value of single biomarkers is never high\n• Depression is gradual, not a sudden episode\n• HRV/sleep are stress markers, not depression-specific", font_size=13, color=LIGHT)
add_box(slide, 6.5, 2.8, 5.3, 2.0, fill_color=BOX_ACCENT_BG)
add_text(slide, 6.8, 2.9, 4.8, 0.3, "Our Answer: Monitor + Intervene", font_size=16, color=ACCENT, bold=True)
add_text(slide, 6.8, 3.2, 4.8, 1.2, "• Detect sleep destabilization early from wearables\n• Sleep is directly treatable (CBT-I, hygiene, meds)\n• Treating sleep improves depression, anxiety, stress", font_size=13, color=LIGHT)
add_box(slide, 0.8, 5.1, 11.5, 0.8)
add_text(slide, 1.1, 5.2, 11, 0.6, "We don't diagnose depression. We flag physiological destabilization — especially sleep —\nso clinicians can intervene before patients drop out.", font_size=14, color=LIGHT, alignment=PP_ALIGN.CENTER)

# ============ SLIDE 10: Two Products ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text(slide, 0.5, 0.3, 12, 0.7, "Two Products, One Engine", font_size=36, color=ACCENT, bold=True)
add_box(slide, 0.8, 1.2, 5.3, 3.0, fill_color=BOX_ACCENT_BG)
add_text(slide, 1.1, 1.3, 4.8, 0.4, "Product 1: Discovery Engine (R&D Tool)", font_size=16, color=ACCENT, bold=True)
add_text(slide, 1.1, 1.7, 4.8, 2.0, "The autoresearch loop — sold to pharma for\ndigital endpoint discovery in clinical trials\n\n• Not a medical device — internal R&D tool\n• No regulatory burden on discovery process\n• Revenue: per-trial licensing to pharma CROs", font_size=13, color=LIGHT)
add_box(slide, 6.5, 1.2, 5.3, 3.0)
add_text(slide, 6.8, 1.3, 4.8, 0.4, "Product 2: Monitoring Dashboard (SaMD)", font_size=16, color=ACCENT2, bold=True)
add_text(slide, 6.8, 1.7, 4.8, 2.0, "Static, validated model for clinical monitoring\n\n• Health Canada SaMD Class II\n• Locked model — no autonomous changes\n• Clinical decision support, not diagnosis\n• Revenue: per-patient SaaS to hospitals", font_size=13, color=LIGHT)
add_box(slide, 0.8, 4.5, 11.5, 1.0)
add_text(slide, 1.1, 4.6, 11, 0.8, "The discovery engine evolves freely. The clinical product is a snapshot —\nfrozen, validated, regulatable. New discoveries trigger a new submission, not a moving target.", font_size=14, color=LIGHT, alignment=PP_ALIGN.CENTER)

# ============ SLIDE 11: Confounder Control ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text(slide, 0.5, 0.3, 12, 0.7, "Confounder-Aware Discovery", font_size=36, color=ACCENT, bold=True)
add_box(slide, 0.8, 1.1, 11.5, 0.7)
add_text(slide, 1.1, 1.2, 11, 0.5, "A \"depression signal\" in HRV might actually be caffeine, alcohol, or exercise.", font_size=17, color=WARN)
add_box(slide, 0.8, 2.2, 5.3, 2.5)
add_text(slide, 1.1, 2.3, 4.8, 0.3, "Confounders in dataset", font_size=16, color=ACCENT2, bold=True)
add_text(slide, 1.1, 2.7, 4.8, 1.6, "Coffee intake        1–5 scale\nAlcohol use           1–4 scale\nSmoking               1–5 scale\nExercise frequency    1–5 scale", font_size=14, color=LIGHT)
add_box(slide, 6.5, 2.2, 5.3, 2.5, fill_color=BOX_ACCENT_BG)
add_text(slide, 6.8, 2.3, 4.8, 0.3, "Our approach", font_size=16, color=ACCENT, bold=True)
add_text(slide, 6.8, 2.7, 4.8, 1.6, "• Residualize features against confounders\n• Report both raw and adjusted effect sizes\n• Flag features where d drops >30% after adjustment\n• FAIR-compliant output: full provenance for each finding", font_size=13, color=LIGHT)
add_text(slide, 0.8, 5.0, 11.5, 0.4, "Every discovery includes: raw effect, adjusted effect, confounders controlled, CI, permutation p-value, FDR status", font_size=11, color=DIM, alignment=PP_ALIGN.CENTER)

# ============ SLIDE 12: What's Next ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text(slide, 0.5, 0.3, 12, 0.7, "What's Next", font_size=36, color=ACCENT, bold=True)
for i, (label, sub) in enumerate([
    ("Hackathon", "Discovery +\nvalidation"),
    ("Clinical\nValidation", "Prospective study\nwith Uni of Hamburg"),
    ("Pharma Pilot", "Discovery engine for\ndigital endpoints"),
    ("SaMD Filing", "Frozen model →\nHealth Canada Class II"),
]):
    x = 0.8 + i * 3.1
    c = BOX_ACCENT_BG if i == 0 else RGBColor(0x1a, 0x15, 0x30)
    add_box(slide, x, 1.3, 2.5, 1.5, fill_color=c)
    add_text(slide, x + 0.1, 1.35, 2.3, 0.5, label, font_size=15, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, x + 0.1, 1.9, 2.3, 0.7, sub, font_size=11, color=DIM, alignment=PP_ALIGN.CENTER)
    if i < 3:
        add_text(slide, x + 2.5, 1.7, 0.6, 0.5, "→", font_size=28, color=ACCENT, alignment=PP_ALIGN.CENTER)
add_box(slide, 0.8, 3.3, 11.5, 1.0, fill_color=BOX_ACCENT_BG)
add_text(slide, 1.1, 3.4, 11, 0.4, "Detect physiological destabilization from a wrist — intervene on sleep before patients drop out.", font_size=17, color=ACCENT, alignment=PP_ALIGN.CENTER)
add_text(slide, 1.1, 3.9, 11, 0.3, "FAIR-compliant outputs • Confounder-controlled • Full provenance trail for every discovery", font_size=12, color=DIM, alignment=PP_ALIGN.CENTER)

# ============ SLIDE 13: Closing ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text(slide, 1, 1.5, 11, 1.0, "AutoBiomarker", font_size=54, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text(slide, 1, 2.7, 11, 0.8, "Autonomous discovery of physiological destabilization\nfrom wearable data", font_size=22, color=ACCENT, alignment=PP_ALIGN.CENTER)
add_box(slide, 2, 3.8, 9, 1.2, fill_color=RGBColor(0x1a, 0x15, 0x30))
add_text(slide, 2.3, 3.9, 8.4, 1.0, "\"Detect sleep and HRV destabilization early.\nTreat sleep — and hit two birds with one stone.\"", font_size=18, color=LIGHT, alignment=PP_ALIGN.CENTER)
add_text(slide, 1, 5.5, 11, 0.4, "Toronto  •  In collaboration with Prof. Dr. Steffen Moritz, Uni of Hamburg", font_size=13, color=DIM, alignment=PP_ALIGN.CENTER)

# Save
out_path = "/Users/athifshaffy/Documents/Repo/COGITO/evolved26-hackathon/AutoBiomarker_Evolved26.pptx"
prs.save(out_path)
print(f"Saved to {out_path}")
